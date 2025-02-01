import os
import re
import streamlit as st
import pyttsx3
import google.generativeai as genai  # Import Gemini
from pptx import Presentation
from io import BytesIO
from langchain_community.document_loaders import PyPDFLoader

# Initialize Gemini 1.5 Flash
genai.configure(api_key="AIzaSyDEH9wwhq4vCiD1lDWLwndGREB7jaeNmfI")  # Replace with your Gemini API key
model = genai.GenerativeModel("gemini-1.5-flash")

# Streamlit UI
st.title("Research Paper Summarizer and PowerPoint Generator with 🎙 Podcast Feature")
st.write("Group Name: Software_Chaser")
st.write("Members: Saiji (Leader), Vyom, Maharsh, Dhwanil, Akshay")
st.write("Upload a PDF file to get a summary, listen to expert discussion audio, and download the PowerPoint presentation.")

# Dropdown for summary level
level_prompts = {
    "Beginner": "Summarize this research paper section for a high school student in 4-6 concise bullet points:",
    "Student": "Create a structured summary of this section for undergraduate students in 4-6 points:",
    "Expert": "Generate a detailed summary of this section for researchers in 5-7 well-formed bullet points:"
}
summary_level = st.selectbox("Select summary level:", list(level_prompts.keys()))

# Dropdown for creativity level
creativity_levels = {
    "Formal": "Keep the conversation strictly professional and formal.",
    "Balanced": "Maintain a balance between professional and conversational tone.",
    "Creative": "Make the conversation more creative and engaging with some informal elements."
}
creativity_level = st.selectbox("Select creativity level:", list(creativity_levels.keys()))

# Dropdown for podcast length
podcast_lengths = {
    "Short (2-3 mins)": "Generate a short podcast with 2-3 questions and concise answers.",
    "Medium (5-7 mins)": "Generate a medium-length podcast with 4-5 questions and detailed answers.",
    "Long (10+ mins)": "Generate a long podcast with 6-8 questions and in-depth discussion."
}
podcast_length = st.selectbox("Select podcast length:", list(podcast_lengths.keys()))

# Dropdown for PowerPoint template selection
template_options = {
    "Template 1": "theme_template_1.pptx",
    "Template 2": "theme_template_2.pptx",
    "Template 3": "theme_template_3.pptx"
}
selected_template = st.selectbox("Select PowerPoint Template:", list(template_options.keys()))

# Custom CSS for avatars
st.markdown("""
<style>
    .avatar-container {
        display: flex;
        justify-content: space-around;
        margin: 2rem 0;
        padding: 20px;
        background: #f0f2f6;
        border-radius: 15px;
    }
    .avatar-card {
        text-align: center;
        transition: all 0.3s ease;
        padding: 15px;
        border-radius: 10px;
    }
    .avatar-img {
        width: 120px;
        height: 120px;
        border-radius: 50%;
        transition: all 0.3s ease;
        filter: grayscale(80%) brightness(0.8);
        border: 3px solid transparent;
    }
    .active-speaker .avatar-img {
        filter: none;
        transform: scale(1.1);
        border-color: #4CAF50;
        box-shadow: 0 0 20px rgba(76, 175, 80, 0.3);
    }
    .avatar-name {
        margin-top: 10px;
        font-weight: bold;
        color: #666;
    }
    .active-speaker .avatar-name {
        color: #4CAF50;
    }
    @keyframes pulse {
        0% { transform: scale(1); }
        50% { transform: scale(1.05); }
        100% { transform: scale(1); }
    }
    .speaking-indicator {
        width: 15px;
        height: 15px;
        background: #4CAF50;
        border-radius: 50%;
        margin: 10px auto;
        opacity: 0;
        animation: pulse 1s infinite;
    }
    .active-speaker .speaking-indicator {
        opacity: 1;
    }
</style>
""", unsafe_allow_html=True)

# File uploader for PDF
pdf_file = st.file_uploader("Upload PDF", type=["pdf"])

def get_avatar_html(active_speaker):
    return f"""
    <div class="avatar-container">
        <div class="avatar-card {'active-speaker' if active_speaker == 'Alex' else ''}">
            <div class="speaking-indicator"></div>
            <img src="https://img.icons8.com/color/144/000000/circled-user-female-skin-type-5.png" 
                 class="avatar-img">
            <div class="avatar-name">Alex</div>
        </div>
        <div class="avatar-card {'active-speaker' if active_speaker == 'Dr. Smith' else ''}">
            <div class="speaking-indicator"></div>
            <img src="https://img.icons8.com/color/144/000000/circled-user-male-skin-type-7.png" 
                 class="avatar-img">
            <div class="avatar-name">Dr. Smith</div>
        </div>
    </div>
    """

def extract_and_summarize_sections(text, summary_level):
    """Extract sections and generate summaries using Gemini."""
    prompt = f"""Analyze the following research paper and:
1. Identify all major sections.
2. For each section, generate a summary using the following guidelines:
   - {level_prompts[summary_level]}
3. Format the response as:
   ## Section Name
   - Bullet point 1
   - Bullet point 2
   - Bullet point 3

Paper content:
{text}
"""
    try:
        response = model.generate_content(prompt)  # Use the initialized model
        return response.text
    except Exception as e:
        st.error(f"Error processing document: {e}")
        return None

def create_ppt_from_summary(summary_text, template_path):
    """Create PowerPoint from section-wise summaries using the selected template."""
    prs = Presentation(template_path)

    # Generate a title for the presentation using Gemini
    title_prompt = f"""Analyze the following text and generate a concise, professional title for a PowerPoint presentation (maximum 10-12 words):
    {summary_text[:5000]}  # Use the first 5000 characters for title generation
    """
    try:
        title_response = model.generate_content(title_prompt)
        title = title_response.text.strip()
    except Exception as e:
        st.error(f"Error generating title: {e}")
        title = "Research Summary"  # Default title if title generation fails

    # Ensure the title fits within one or two lines
    max_title_length = 80  # Maximum characters per line (adjust as needed)
    if len(title) > max_title_length:
        # Split the title into two lines
        words = title.split()
        line1 = ""
        line2 = ""
        for word in words:
            if len(line1) + len(word) + 1 <= max_title_length:  # +1 for space
                line1 += word + " "
            else:
                line2 += word + " "
        title = f"{line1.strip()}\n{line2.strip()}"

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]  # 0 is the layout for a title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by AI Podcast Generator"

    # Add section slides
    slide_layout = prs.slide_layouts[1]  # 1 is the layout for a content slide

    # Parse the summary text into sections and bullet points
    sections = {}
    current_section = "Introduction"
    for line in summary_text.split('\n'):
        if line.startswith("## "):
            current_section = line[3:].strip()
            sections[current_section] = []
        elif line.startswith("- "):
            sections[current_section].append(line[2:])

    # Create slides for each section
    for section, bullets in sections.items():
        slides_per_section = min((len(bullets) // 6) + 1, 5)
        chunk_size = max(len(bullets) // slides_per_section, 1)
        
        for i in range(0, len(bullets), chunk_size):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = section if i == 0 else f"{section} (Cont.)"
            
            content_box = slide.shapes.placeholders[1]
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for bullet in bullets[i:i+chunk_size]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = 0
    
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

def generate_podcast_script(summary_text, creativity_level, podcast_length):
    """Generate a conversational podcast script using Gemini."""
    prompt = f"""Create a conversational podcast script between host Alex and researcher Dr. Smith discussing the research paper. Follow these rules:
1. Alex should ask curious, layperson-friendly questions.
2. Dr. Smith should provide expert answers based on the paper.
3. Always prefix lines with either "Alex:" or "Dr. Smith:".
4. Keep responses conversational but informative.
5. Cover key findings, methodology, and implications.
6. {creativity_levels[creativity_level]}
7. {podcast_lengths[podcast_length]}

Paper summary:
{summary_text}
"""
    try:
        response = model.generate_content(prompt)  # Use the initialized model
        return response.text
    except Exception as e:
        st.error(f"Error generating podcast script: {e}")
        return None

def generate_podcast_audio(podcast_script, rate=150):
    """Generate TTS audio with distinct voices for host and researcher."""
    try:
        engine = pyttsx3.init()
        engine.setProperty('rate', rate)
        voices = engine.getProperty('voices')
        male_voice = voices[0].id
        female_voice = voices[1].id
        avatar_placeholder = st.empty()
        for line in podcast_script.split('\n'):
            line = line.strip()
            if not line:
                continue

            if line.startswith("Alex:"):
                engine.setProperty('voice', female_voice)
                text = line.replace("Alex:", "").strip()
                avatar_placeholder.markdown(get_avatar_html('Alex'), unsafe_allow_html=True)
            elif line.startswith("Dr. Smith:"):
                engine.setProperty('voice', male_voice)
                text = line.replace("Dr. Smith:", "").strip()
                avatar_placeholder.markdown(get_avatar_html('Dr. Smith'), unsafe_allow_html=True)
            else:
                continue

            engine.say(text)
            engine.runAndWait()

        return True
    except Exception as e:
        st.error(f"Error generating audio: {e}")
        return False

if pdf_file is not None:
    # Save the uploaded PDF temporarily
    temp_pdf_path = f"./temp_{pdf_file.name}"
    with open(temp_pdf_path, "wb") as f:
        f.write(pdf_file.getbuffer())

    # Load and process the PDF
    loader = PyPDFLoader(temp_pdf_path)
    documents = loader.load()
    pdf_text = "\n".join([doc.page_content for doc in documents])

    # Extract sections and generate summaries using Gemini
    with st.spinner("Analyzing document and generating summaries..."):
        summary_text = extract_and_summarize_sections(pdf_text, summary_level)

    if summary_text:
        # Display summaries
        st.subheader("Section-wise Summary")
        st.markdown(summary_text)

        # Generate podcast content
        podcast_script = generate_podcast_script(summary_text, creativity_level, podcast_length)

        if podcast_script:
            st.subheader("Expert Discussion Script")
            st.text_area(
            "Expert Discussion Script",
            value=podcast_script,
            height=300,  # Set the height of the textarea
            key="podcast_script_area"
             )
            
            # Play podcast audio
            st.subheader("Listen to Expert Discussion")
            if st.button("Play Podcast"):
                if generate_podcast_audio(podcast_script):
                    st.success("Podcast audio played successfully!")
                else:
                    st.error("Failed to play podcast audio.")

        # Generate and download PowerPoint
        pptx_stream = create_ppt_from_summary(summary_text, template_options[selected_template])
        
        st.subheader("Download Presentation Slides")
        st.download_button(
            label="Download PowerPoint",
            data=pptx_stream,
            file_name="research_summary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    # Cleanup
    os.remove(temp_pdf_path)